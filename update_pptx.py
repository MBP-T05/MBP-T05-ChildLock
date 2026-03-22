import collections
import collections.abc
import os
try:
    import pptx
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx is not installed.")
    exit(1)

def replace_text(shape, replacements):
    if not shape.has_text_frame:
        return
    for prg in shape.text_frame.paragraphs:
        for run in prg.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def update_presentation():
    base_dir = r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock"
    pptx_path = os.path.join(base_dir, r"docs\MBP-T05_ChildLock_결과발표.pptx")
    out_path = os.path.join(base_dir, r"docs\MBP-T05_ChildLock_결과발표_Updated.pptx")
    
    prs = pptx.Presentation(pptx_path)

    # --- Slide 1 (Index 0) ---
    s1_reps = {
        "ISO 26262 SW 시뮬레이션 프로젝트": "ISO 26262 기반 소프트웨어 V-모델 V&V 적용",
        "자율주행 제어 및 검증 기반": "ISO 26262 및 A-SPICE 기반 프로세스",
        "전차": ""
    }
    for shape in prs.slides[0].shapes:
        replace_text(shape, s1_reps)

    # --- Slide 3 (Index 2) ---
    s3_reps = {
        "문제점 1: 기계식 차일드 락의 운전자 직접 제어 불편함": "문제점 1: 기계식 차일드 락의 운전자 직접 제어 불편함 및 확인 어려움\n문제점 2: 보호자가 뒷좌석 상황을 인지하지 못한 상태에서 문 열림 시 사고 위험",
        "목표 1: 속도 기반 자동 잠금으로 편의성 개선": "목표 1: 뒷좌석 통합 전자 제어(운전석 제어) 도입을 통한 편의성 개선\n목표 2: 후방 레이더 연동 스마트 하차 안전 보조(FG-02) 기능 도입",
        "ISO 26262를 준수하면서 승객의 편의성과 하차 안전을 모두 만족하는 지능형 차일드 락 제어기를 어떻게 구현할 것인가?": "ISO 26262 기능 안전을 준수하면서, 승객의 편의성과 하차 안전을 모두 충족하는 지능형 시스템을 어떻게 구현·검증할 것인가?"
    }
    for shape in prs.slides[2].shapes:
        replace_text(shape, s3_reps)
        
    # --- Slide 4 (Index 3) ---
    s4_reps = {
        "C/C++ 기반 ISO 26262 안전 요구사항 적용 제어 로직 모듈 탑재": "C/C++ 기반 ISO 26262 및 FMEA/HAZOP 상위 기능 안전 분석 적용 모듈",
    }
    for shape in prs.slides[3].shapes:
        replace_text(shape, s4_reps)
        
    s4_shapes = [s for s in prs.slides[3].shapes if s.has_text_frame and "속도 센서 연동, 후측방 레이더 결합 제어 로직" in s.text]
    s4_shapes.sort(key=lambda s: s.top)
    if len(s4_shapes) >= 1:
        s4_shapes[0].text_frame.text = "속도, 충돌, 스위치 연동 로직 적용 및 직전 상태 복구(Safe State) 제어"
    if len(s4_shapes) >= 2:
        s4_shapes[1].text_frame.text = "후방 위험 판단 시 자동 잠금 유지 및 뒷좌석 점유/출발 상태 전환 시 운전자 알림 연동"
    if len(s4_shapes) >= 3:
        s4_shapes[2].text_frame.text = "FMEA/HAZOP 기반 요구사항 분석 적용 및 Fault Injection 시나리오 방어를 통한 Fail-Safe 설계"

    # --- Slide 5 (Index 4) ---
    s5_reps = {
        "ISO 26262, MISRA C:2012, ASPICE, CMake": "ISO 26262, ASPICE, V-Model 워크플로우 적용",
        "Google Test v1.17.0, Cppcheck": "Google Test, Cppcheck (단위 테스트 및 정적 분석)",
        "Git, GitHub Actions (CI), Gcov/Lcov": "GitHub Actions (CI 자동화), Draw.io (UML)"
    }
    for shape in prs.slides[4].shapes:
        replace_text(shape, s5_reps)
        
    s5_img_path = os.path.join(base_dir, r"docs\diagram\state\childlock_state.png")
    if os.path.exists(s5_img_path):
        prs.slides[4].shapes.add_picture(s5_img_path, Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))

    # --- Slide 6 (Index 5) ---
    # Find the title shape to position below it.
    slide6 = prs.slides[5]
    s6_img1 = os.path.join(base_dir, r"docs\diagram\flow_chart\FG-01_door_control.png")
    s6_img2 = os.path.join(base_dir, r"docs\diagram\flow_chart\FG-02_exit_safety.png")
    
    if os.path.exists(s6_img1):
        slide6.shapes.add_picture(s6_img1, Inches(0.5), Inches(1.5), Inches(4.0), Inches(4.5))
        tx1 = slide6.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(4.0), Inches(1))
        tx1.text_frame.text = "안전 상태(SAFE_LOCKED) 전이 및 직전 상태 복구 제어 로직"
        tx1.text_frame.paragraphs[0].font.size = Pt(14)
        
    if os.path.exists(s6_img2):
        slide6.shapes.add_picture(s6_img2, Inches(5.0), Inches(1.5), Inches(4.0), Inches(4.5))
        tx2 = slide6.shapes.add_textbox(Inches(5.0), Inches(6.1), Inches(4.0), Inches(1))
        tx2.text_frame.text = "AI 기반 하차 안전 보조(후방 위험 객체 탐지 및 상태 알림 연동)"
        tx2.text_frame.paragraphs[0].font.size = Pt(14)

    # --- Slide 7 (Index 6) ---
    s7_reps = {
        "테스트 항목 (TC-CHL-001 ~ 015 & GTest Suites)": "테스트 항목 (단위 테스트, Fault Injection 시나리오 및 안전 상태 전이 검증)",
        "ASIL B / QM 검증 요구사항": "ASIL B / QM 검증 요구사항 충족 및 양방향 추적성(Traceability) 점검 완료",
    }
    for shape in prs.slides[6].shapes:
        replace_text(shape, s7_reps)

    # --- Slide 8 (Index 7) ---
    s8_shapes = [s for s in prs.slides[7].shapes if s.has_text_frame and "ISO 26262/ASPICE 기능안전 기반 요구사항(SRS)" in s.text]
    s8_shapes.sort(key=lambda s: s.top)
    if len(s8_shapes) >= 1:
        s8_shapes[0].text_frame.text = "ISO 26262/ASPICE 산출물 기반 요구사항(SRS, SwDD) 작성 체화 및 V-모델 기반 설계-구현 워크플로우 경험"
    if len(s8_shapes) >= 2:
        s8_shapes[1].text_frame.text = "GitHub Action을 활용한 CI/CT 자동 검증 파이프라인 연동 및 정적 분석 룰셋(Cppcheck) 도입 경험"

    s8_reps2 = {
        "복합 센서 예외 조건 처리 로직의 복잡성 증가 및 Fault 발생 시나리오 정리 애로": "정상 동작 외의 복잡한 예외 조건(센서 에러, 동시 입력 등) 처리 시 상태 전이 충돌 위험 증가",
        "Fault Injection 테스트 케이스와 Fail-Safe 시나리오(Decision Table) 구체화 적용": "제어 우선순위를 반영한 Decision Table 설계로 논리적 결함 차단 및 Fault Injection 시나리오화로 해결"
    }
    for shape in prs.slides[7].shapes:
        replace_text(shape, s8_reps2)

    # --- Slide 9 (Index 8) ---
    s9_shapes = [s for s in prs.slides[8].shapes if s.has_text_frame and "제어기 로직 고도화 적용" in s.text]
    s9_shapes.sort(key=lambda s: s.top)
    if len(s9_shapes) >= 1:
        s9_shapes[0].text_frame.text = "다중 제어기 통합 시그널 최적화 연동"
    if len(s9_shapes) >= 2:
        s9_shapes[1].text_frame.text = "결합 시나리오 확장 및 판단 신뢰성 증대"

    try:
        prs.save(out_path)
        print(f"Presentation successfully saved to {out_path}")
    except Exception as e:
        print(f"Failed to save presentation: {e}")

if __name__ == "__main__":
    update_presentation()
