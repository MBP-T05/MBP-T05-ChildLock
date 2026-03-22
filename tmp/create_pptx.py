from pptx import Presentation

def main():
    prs = Presentation('docs/MOBIUS_PBL_결과발표_템플릿.pptx')
    
    # ---------------------------------------------------------
    # Slide 1 (Index 0): Title Slide
    # Pld 0: Title 1 -> 프로젝트/팀명
    # Pld 1: Subtitle 2 -> 추가 정보
    # Pld 10: Date Placeholder 3 -> 2026. 03.
    # ---------------------------------------------------------
    s1 = prs.slides[0]
    for shape in s1.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = "전자식 차일드 락 시스템 (MBP-T05)\nISO 26262 SW 시뮬레이션 프로젝트"
        elif shape.placeholder_format.idx == 1:
            shape.text = "팀 전차\n박기준(팀장), 김도균, 김이안, 이한결, 이승욱, 박찬석"
        elif shape.placeholder_format.idx == 10:
            shape.text = "2026. 03. 22"

    # ---------------------------------------------------------
    # Slide 2 (Index 1): 01. 팀 소개
    # Pld 0: Title
    # Pld 1: Content
    # ---------------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = "박기준 (팀장)\n• 프로젝트 총괄, 시스템 아키텍처 및 요구사항 정의\n• ISO 26262 안전 메커니즘 설계\n\n팀원 (김도균, 김이안, 이한결, 이승욱, 박찬석)\n• 도어 및 하차 안전 제어 로직 기획 및 구현\n• 단위 테스트(GTest), 정적 분석(Cppcheck) 등 테스트 활동\n• 지속적 통합(CI) 환경 구축 및 빌드 자동화"

    # ---------------------------------------------------------
    # Slide 3 (Index 2): 02. 문제 정의
    # Pld 1: AS-IS / TO-BE
    # ---------------------------------------------------------
    s3 = prs.slides[2]
    # Finding text placeholders based on previously inspected text
    for shape in s3.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "AS-IS" in text:
             shape.text = "AS-IS (현재 문제점)\n• 문제점 1: 수동 기계식 차일드 락의 조작 불편 (운전자 직접 변경 필요)\n• 문제점 2: 사고 시 뒷좌석 승객 탈출 지연 위험\n• 문제점 3: 후방 위험(자전거/오토바이 등) 인지 및 제어 수단 부재"
        elif "TO-BE" in text:
             shape.text = "TO-BE (목표 상태)\n• 목표 1: 속도 기반 자동 잠금 제어를 통한 편의 개선\n• 목표 2: 충돌 등 비상 상황 시 뒷좌석 자동 잠금 해제로 대피 지원\n• 목표 3: 후방 센서 연동을 통한 하차 시 보호 로직 탑재"
        elif "핵심 연구 질문" in text:
             shape.text = "핵심 연구 질문\nISO 26262 안전 표준을 준수하면서 승객의 편의성과 하차 안전을 모두 만족하는 지능형 차일드 락 SW를 어떻게 구현할 것인가?"

    # ---------------------------------------------------------
    # Slide 4 (Index 3): 03. 솔루션 개요
    # ---------------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "솔루션 한 줄 요약" in text:
            shape.text = "지능형 센서 및 차량 속도 연동 안전 제어 소프트웨어"
        elif "핵심 기능 1" in text:
            shape.text = "FG-01 도어 잠금/해제 제어\n• 운전자 스위치 조작 제어\n• 차량 속도 연동 자동 잠금 활성화"
        elif "핵심 기능 2" in text:
            shape.text = "비상 상황 대처\n• 차량 충돌 발생 시 뒷좌석 비상 잠금 해제 제어\n• 상시 실내 탑승자 개방 금지 메커니즘 결합"
        elif "핵심 기능 3" in text:
            shape.text = "FG-02 하차 안전 보조 제어\n• 후측방 차량/이륜차 접근 감지 연동\n• 뒷좌석 점유 및 상태 요약 알림"

    # ---------------------------------------------------------
    # Slide 5 (Index 4): 04. 기술 스택 & 아키텍처
    # ---------------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "다이어그램" in text:
            shape.text = "[시스템 아키텍처 다이어그램 공간]\n• State Machine: Locked, Unlocked, Emergency 상태 전이\n• Logic Flow: 센서 입력 -> 제어 판별 -> Actuator 제어"
        elif "Tech Stack" in text:
            shape.text = "Tech Stack\nLanguage\nC11 / C++17\nFramework\nISO 26262, MISRA C:2012, ASPICE\nTesting\nGoogle Test v1.17.0, Cppcheck\nTools\nCMake, Lizard, Gcov / Lcov\nOS\nUbuntu 22.04 LTS"

    # ---------------------------------------------------------
    # Slide 6 (Index 5): 05. 구현 결과
    # ---------------------------------------------------------
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "스크린샷 삽입" in text:
            if "1" in text:
                shape.text = "[유스케이스 1]\n속도(3km/h) 초과 시 자동 잠금"
            elif "2" in text:
                shape.text = "[유스케이스 2]\n충돌 감지 시 즉시 잠금 해제"
            elif "3" in text:
                shape.text = "[유스케이스 3]\n후측방 위험 물체 접근 시 문열림 차단"
            elif "4" in text:
                shape.text = "[유스케이스 4]\n시동 종료 전 차일드락 상태 요약 표시"

    # ---------------------------------------------------------
    # Slide 7 (Index 6): 06. 테스트 & 검증
    # ---------------------------------------------------------
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "테스트 케이스" in text:
            shape.text = "테스트 케이스 및 결과 (단위 및 서브시스템 검증)\n• UC-01~07 전 구간 GTest 기반 테스트 작성\n• 상태 머신 (State Machine) 전이 조건 만족 검증"
        if "성과 지표" in text:
            pass # Keep table headers
        if "정확도" in text:
            shape.text = "성과 지표\n정확도\n100%\n처리 속도\n<1ms\n테스트 통과율\n100%\n커버리지\n100%"

    # ---------------------------------------------------------
    # Slide 8 (Index 7): 07. 학습 내용 & 교훈
    # ---------------------------------------------------------
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "기술적 학습" in text:
            shape.text = "기술적 학습\n• ISO 26262 기반 요구사항 도출 및 시스템 설계 방법론 체득\n• MISRA C 코딩 가이드라인 실무 적용\n• Cppcheck 및 GTest를 활용한 엄격한 품질보증 파이프라인"
        if "협업 & 프로세스" in text:
            shape.text = "협업 & 프로세스\n• CI/CD 파이프라인(CMake + GTest)을 통한 테스트 자동화 구축\n• GitHub Policy 및 커밋 컨벤션 준수 중요성"
        if "도전 과제" in text:
            shape.text = "도전 과제 & 해결 과정\n[도전 과제]\n다양한 센서 이벤트 및 예외 상황 엣지케이스 테스트 난이도\n[해결 방법]\nTDD 방법론 도입 및 Failure Mode (FMEA/HARA) 사전 정의 후 분기 모의"

    # ---------------------------------------------------------
    # Slide 9 (Index 8): 08. 향후 계획
    # ---------------------------------------------------------
    s9 = prs.slides[8]
    for shape in s9.shapes:
        if not shape.has_text_frame: continue
        text = shape.text
        if "1단계" in text:
            shape.text = "1단계\n[통신 연동 강화]\n• CAN(차량 내 네트워크) 통신 레이어 모의 통합 반영\n• 다양한 ECU 신호 통합 시뮬레이션"
        elif "2단계" in text:
            shape.text = "2단계\n[하드웨어 포팅]\n• 실제 모터 제어(Actuator) 보드 포팅 및 검증\n• 실시간 운영체제(RTOS) 환경 최적화"
        elif "3단계" in text:
            shape.text = "3단계\n[실용화]\n• 고도화된 타겟 환경 통합 플랫폼 시연\n• 동적 분석 추가 도입을 통한 기능 안전 규격 확보 강화"
        elif "실무 적용 가능성" in text:
            shape.text = "실무 적용 가능성\n• 현대모비스/협력사의 차세대 통합 차량 제어기 플랫폼\n• 하차 안전 보조와 차일드 락 기능 간소화 및 통합 지능화 모듈 솔루션"
        elif "추가 개발 아이디어" in text:
            shape.text = "추가 개발 아이디어\n• 모바일 앱 제어 API 제공 (블루투스 연동)\n• 레이더/라이다 기반 더욱 정밀한 후측방 접근 객체 분류"

    prs.save('docs/MBP-T05_ChildLock_결과발표.pptx')
    print("성공적으로 docs/MBP-T05_ChildLock_결과발표.pptx 파일이 생성되었습니다.")

if __name__ == "__main__":
    main()
