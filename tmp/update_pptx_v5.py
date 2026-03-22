"""
MBP-T05 Child Lock - PPT 전면 개선 스크립트 (학회 발표 스타일)
Output: docs/MBP-T05_ChildLock_결과발표_Final_v2.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── 색상 팔레트 ─────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)

BLUE_DARK  = rgb(0x00, 0x35, 0x7e)
BLUE_MID   = rgb(0x00, 0x5b, 0xb5)
BLUE_LIGHT = rgb(0xe8, 0xf0, 0xfe)
ACCENT     = rgb(0x00, 0xa8, 0xe8)
WHITE      = rgb(0xff, 0xff, 0xff)
DARK_TEXT  = rgb(0x1a, 0x1a, 0x2e)
GRAY_BG    = rgb(0xf4, 0xf6, 0xf9)
GREEN      = rgb(0x00, 0x96, 0x60)
ASIL_B     = rgb(0xff, 0x6b, 0x35)
LITE_BLUE  = rgb(0xcc, 0xe5, 0xff)

def hex6(color):
    """RGBColor to 6-char hex string"""
    return '%02X%02X%02X' % (color[0], color[1], color[2])

BASE = r'c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock'
DIAG = os.path.join(BASE, 'docs', 'diagram')
SRC  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Final (1).pptx')
OUT  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Final_v2.pptx')

# ─── 유틸리티 ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """단색 사각형 배경 추가"""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 선 없음
    return shape


def add_tb(slide, left, top, width, height, text,
           size=11, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트 박스 추가 (배경 투명)"""
    if color is None:
        color = DARK_TEXT
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Malgun Gothic'
    return txb


def add_tb_multiline(slide, left, top, width, height, lines,
                     size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트 박스"""
    if color is None:
        color = DARK_TEXT
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size if line_size else size)
        run.font.bold = line_bold if line_bold is not None else bold
        run.font.color.rgb = color
        run.font.name = 'Malgun Gothic'
    return txb


def fit_image(slide, img_path, left, top, max_w, max_h):
    """이미지를 비율 유지하며 영역 내 중앙 배치"""
    try:
        from PIL import Image
        img = Image.open(img_path)
        iw, ih = img.size
        ratio = min(max_w / iw, max_h / ih)
        nw = iw * ratio
        nh = ih * ratio
        cx = left + (max_w - nw) / 2
        cy = top + (max_h - nh) / 2
        return slide.shapes.add_picture(
            img_path, Inches(cx), Inches(cy), Inches(nw), Inches(nh))
    except ImportError:
        # PIL 없으면 그냥 width로 맞춤
        try:
            return slide.shapes.add_picture(
                img_path, Inches(left), Inches(top), Inches(max_w))
        except Exception as e:
            print(f'  [WARN] Image {os.path.basename(img_path)}: {e}')
            return None
    except Exception as e:
        print(f'  [WARN] Image {os.path.basename(img_path)}: {e}')
        return None

def remove_shapes_with_text(slide, keywords):
    """특정 키워드를 포함한 shape 제거"""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text.strip()
            if any(kw in t for kw in keywords):
                to_remove.append(shape)
        # Check tables too
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if any(kw in t for kw in keywords):
                        to_remove.append(shape)
                        break
    for s in to_remove:
        try:
            s.element.getparent().remove(s.element)
        except Exception:
            pass


def add_footer(slide, page_num, total=10):
    add_rect(slide, 0, 6.82, 13.33, 0.32, BLUE_DARK)
    add_tb(slide, 0.1, 6.85, 13.0, 0.28,
           f'MOBIUS BOOTCAMP 1기 | PBL 결과발표 | {page_num}/{total}',
           size=8, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Slide 2: 팀원 카드 재구성 ───────────────────────────────────────────────

def rebuild_slide2(slide):
    print('  Slide 2: 팀원 카드 재구성')
    remove_shapes_with_text(slide, [
        '박기준', '김도균', '김이안', '이한결', '이승욱', '박찬석',
        '프로젝트 리더', '담당 업무', '역할', 'Core Developer',
        '도어 및 하차', '단위 테스트'
    ])

    members = [
        ('이한결', 'hangyeoli', '프로젝트 메인테이너 & F-01',
         'F-01 InputMonitorAndValidator',
         ['전체 Issues 생성 관리 (Issue #4~#13)',
          'CMake 빌드 구성 (PR#17)', 'Scope Fix & CTest (PR#26)'],
         BLUE_DARK),
        ('김이안', 'iank1m', '상태 결정 로직 개발자',
         'F-02 ChildLockStateDecision',
         ['차일드락 핵심 상태 전이 로직',
          '단독 구현 (PR#14)',
          'ISO 26262 ASIL 설계 반영'],
         BLUE_MID),
        ('박찬석', 'p-chanseok', 'TDD 기반 제어 로직 개발자',
         'F-03 DoorEcuCommandHandler\nF-04 RearDoorOpenBlockHandler',
         ['TDD 방식 Door ECU 제어 (PR#23)',
          '69/69 테스트 케이스 통과 (PR#22)'],
         BLUE_DARK),
        ('이승욱', 'josephuk77', '상태 유지 & 알림 개발자',
         'F-05 StatePersistenceManager\nF-10 IgnitionOffStatusAlert',
         ['상태 저장/복구 및 시동 OFF 알림 (PR#19, #20)',
          'UML 다이어그램 문서화 (Issue#1/PR#2)'],
         BLUE_MID),
        ('김도균', 'codingSkeleton6478', 'HMI & 좌석 알림 개발자',
         'F-06 HmiAndEventLogger\nF-09 RearSeatOccupancyAlert',
         ['경고/알림 HMI 출력 로직 (PR#24)',
          '이벤트 로거 구현'],
         BLUE_DARK),
        ('박기준 (팀장)', 'PKJthecreator', 'PM & Safety 모듈 개발자',
         'F-07 RearRiskEvaluation\nF-08 RearRiskProtectionController',
         ['후방 위험 평가/보호 핵심 ASIL B 로직',
          '(PR#15, PR#16) + PM 총괄'],
         BLUE_MID),
    ]

    # 3×2 카드 레이아웃
    card_w, card_h = 4.1, 2.55
    gap = 0.13
    col_x = [0.13, col_x1 := 0.13 + card_w + gap, 0.13 + 2*(card_w+gap)]
    row_y = [1.32, 1.32 + card_h + 0.12]
    col_positions = [0.13, 0.13 + card_w + gap, 0.13 + 2*(card_w+gap)]

    for i, m in enumerate(members):
        col = i % 3
        row = i // 3
        x = col_positions[col]
        y = row_y[row]
        name, gid, role, module, details, hdr_color = m

        # 헤더
        add_rect(slide, x, y, card_w, 0.8, hdr_color)
        add_tb(slide, x+0.08, y+0.06, card_w-0.16, 0.38,
               name, size=13, bold=True, color=WHITE)
        add_tb(slide, x+0.08, y+0.44, card_w-0.16, 0.32,
               f'@{gid}  |  {role}',
               size=8.5, bold=False, color=LITE_BLUE)

        # 본문 배경
        add_rect(slide, x, y+0.8, card_w, card_h-0.8, GRAY_BG)

        # 모듈명
        add_tb(slide, x+0.1, y+0.85, card_w-0.2, 0.45,
               module, size=9, bold=True, color=hdr_color)

        # 상세 내용
        detail_y = y + 1.35
        for line in details:
            add_tb(slide, x+0.12, detail_y, card_w-0.24, 0.28,
                   f'• {line}', size=8, color=DARK_TEXT)
            detail_y += 0.28


# ─── Slide 7: 테스트 케이스 표 (index 6, "06. 테스트 케이스 및 결과") ──────────

def rebuild_slide7(slide):
    print('  Slide 7 (Index 6): 테스트 케이스 표 재구성')
    remove_shapes_with_text(slide, [
        '정상 동작 검증', '테스트 항목', '성과 지표',
        'Pass/Fail', '정확도', '처리 속도', '커버리지', '테스트 통과율',
        'MOBIUS BOOTCAMP'
    ])
    # 테이블 shape 제거
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            to_remove.append(shape)
    for s in to_remove:
        try:
            s.element.getparent().remove(s.element)
        except Exception:
            pass

    # TC 데이터: ReqTest.md 반영
    tcs = [
        ('TC-CHL-001', 'UC-1', '정차 중 수동 스위치 제어 정상 동작', 'EP', 'QM'),
        ('TC-CHL-002', 'UC-1', '스위치 채터링 디바운싱 - 49ms (Off-limit)', 'BVA', 'QM'),
        ('TC-CHL-003', 'UC-1', '스위치 채터링 디바운싱 - 50ms (On-limit)', 'BVA', 'QM'),
        ('TC-CHL-004', 'UC-1', '주행 중 해제 차단 - 3.0km/h (On-limit)', 'BVA', 'ASIL B'),
        ('TC-CHL-005', 'UC-1', '주행 중 해제 허용 - 2.9km/h (Off-limit)', 'BVA', 'ASIL B'),
        ('TC-CHL-006', 'UC-1', 'ECU 통신 단절 재전송 및 Fail-Safe 안전상태', 'FI', 'QM'),
        ('TC-CHL-007', 'UC-2', '속도 초과 자동잠금 미동작 - 14.9km/h', 'BVA', 'QM'),
        ('TC-CHL-008', 'UC-2', '속도 초과 자동잠금 동작 - 15.0km/h', 'BVA', 'QM'),
        ('TC-CHL-009', 'UC-3', '충돌 시 비상 해제 (복합조건 = True)', 'DT', 'ASIL B'),
        ('TC-CHL-010', 'UC-3', '충돌 단일조건 무시 (오동작 방지)', 'DT', 'ASIL B'),
        ('TC-CHL-011', 'UC-4', 'ECU 상태 전이(Init→Normal→Lock) 핸들 차단', 'ST', 'QM'),
        ('TC-CHL-012', 'UC-5', '후방 객체(거리/속도/시간) 복합 감지', 'PW', 'QM'),
        ('TC-CHL-013', 'UC-5', '후방 센서 노이즈/반복 경고 히스테리시스', 'EP', 'QM'),
        ('TC-CHL-014', 'UC-6', '승객 점유 및 잠금 OFF 출발 시 알림', 'DT', 'QM'),
        ('TC-CHL-015', 'UC-6', '점유 센서 에러 발생 시 보호 로직', 'FI', 'QM'),
    ]
    hdrs = ['Test ID', 'UC', '테스트 항목', '기법', 'ASIL']
    cws  = [0.95, 0.42, 4.05, 0.55, 0.72]
    total_w = sum(cws)
    sx = (13.33 - total_w) / 2
    sy = 1.32
    hh = 0.36
    rh = 0.285

    # 헤더 행
    xp = sx
    for h, cw in zip(hdrs, cws):
        add_rect(slide, xp, sy, cw, hh, BLUE_DARK)
        add_tb(slide, xp+0.03, sy+0.07, cw-0.06, hh-0.1,
               h, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        xp += cw

    # 데이터 행
    for ri, tc in enumerate(tcs):
        ry = sy + hh + ri * rh
        is_asil = tc[4] == 'ASIL B'
        row_bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
        xp = sx
        for ci, (val, cw) in enumerate(zip(tc, cws)):
            cbg = row_bg
            txt_color = DARK_TEXT
            bld = False
            if ci == 4 and is_asil:
                cbg = ASIL_B
                txt_color = WHITE
                bld = True
            add_rect(slide, xp, ry, cw, rh, cbg)
            add_tb(slide, xp+0.03, ry+0.04, cw-0.06, rh-0.06,
                   val, size=8, bold=bld, color=txt_color,
                   align=PP_ALIGN.CENTER)
            xp += cw

    # KPI 박스 (CodeTest.md 연동)
    kpi_y = sy + hh + len(tcs) * rh + 0.12
    kpis = [
        ('단위 테스트', '140 건', BLUE_DARK),
        ('통과율', '100%  (0 Fail)', GREEN),
        ('커버리지', '100%  Branch', BLUE_MID),
        ('처리 속도', '< 1 ms', ACCENT),
        ('시스템 TC', '15 건', ASIL_B),
        ('프레임워크', 'GTest v1.17', BLUE_DARK),
    ]
    kw = 2.1
    kg = 0.1
    total_kw = len(kpis)*(kw+kg) - kg
    ksx = (13.33 - total_kw) / 2
    kh = 0.82

    for ki, (lbl, val, kc) in enumerate(kpis):
        kx = ksx + ki*(kw+kg)
        add_rect(slide, kx, kpi_y, kw, 0.28, kc)
        add_tb(slide, kx+0.04, kpi_y+0.04, kw-0.08, 0.22,
               lbl, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, kx, kpi_y+0.28, kw, kh-0.28, rgb(0xf9, 0xfb, 0xff))
        add_tb(slide, kx+0.04, kpi_y+0.3, kw-0.08, kh-0.32,
               val, size=12, bold=True, color=kc, align=PP_ALIGN.CENTER)

    add_footer(slide, 7)

# ─── Slide 8: 07 학습내용& 교훈 (index 7) ──────────────────────────────────
def rebuild_slide8(slide):
    print('  Slide 8 (Index 7): 학습내용 및 교훈 삽입')
    remove_shapes_with_text(slide, ['학습내용', '교훈', '결론', 'MOBIUS', 'TDD', 'CI', '빌드', '품질'])

    # 기존 제목 살리거나 투명 배경 추가 등을 위한 공간 (1.3 높이부터 사용)
    # 좌측 영역: TDD 기반 검증 및 커버리지 확보
    add_rect(slide, 0.8, 1.8, 5.6, 0.5, BLUE_DARK)
    add_tb(slide, 0.9, 1.9, 5.4, 0.3,
           '🔷 TDD 기반 검증 및 커버리지 확보', size=14, bold=True, color=WHITE)
    
    add_rect(slide, 0.8, 2.3, 5.6, 3.5, GRAY_BG)
    add_tb_multiline(slide, 1.0, 2.5, 5.2, 3.0, [
        ('• GTest 프레임워크를 활용하여 각 컴포넌트 단위 TDD 환경 구축', False, 12),
        ('• 극단적인 엣지 케이스 및 경계값 등 다양한 예외 상황 커버리지 증가', False, 12),
        ('• 결함 주입(Fault Injection) 테스트를 통해 방어적 프로그래밍 정착', False, 12),
        ('• 140건의 단위 테스트 통과 및 Branch/Statement 커버리지 100% 모의 달성', False, 12)
    ], size=13, color=DARK_TEXT)

    # 우측 영역: 자동화 기반 품질 검증 및 MISRA 준수
    add_rect(slide, 6.9, 1.8, 5.6, 0.5, BLUE_MID)
    add_tb(slide, 7.0, 1.9, 5.4, 0.3,
           '🔷 자동화 기반 품질 검증 및 MISRA 준수', size=14, bold=True, color=WHITE)
    
    add_rect(slide, 6.9, 2.3, 5.6, 3.5, GRAY_BG)
    add_tb_multiline(slide, 7.1, 2.5, 5.2, 1.0, [
        ('• GitHub Actions 기반 CI (지속적 통합) 파이프라인 구축', False, 12),
        ('• PR 생성 단계에서 빌드, 테스트, 커버리지, 정적 분석을 자동 수행', False, 12)
    ], size=13, color=DARK_TEXT)

    # 툴 박스들 (하단)
    add_rect(slide, 7.2, 3.6, 5.0, 1.8, WHITE)
    tools = [
        ('Static Analysis', 'Cppcheck'),
        ('Complexity', 'Lizard (≤10)'),
        ('Coverage', 'lcov'),
        ('Code Metrics', 'cloc')
    ]
    ty = 3.8
    for lbl, desc in tools:
        add_tb(slide, 7.4, ty, 2.0, 0.3, f'[{lbl}]', size=11, bold=True, color=BLUE_DARK)
        add_tb(slide, 9.4, ty, 2.5, 0.3, desc, size=11, bold=True, color=DARK_TEXT)
        ty += 0.4

    add_footer(slide, 8)


# ─── Slide 10: 팀명 교체 ─────────────────────────────────────────────────────

def fix_slide10(slide):
    print('  Slide 10 (팀명 확인 및 교체)')
    for shape in slide.shapes:
        if shape.has_text_frame and '[팀명]' in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '[팀명]' in run.text:
                        run.text = run.text.replace('[팀명]', '전차 팀')


def main():
    print(f'Loading: {SRC}')
    prs = Presentation(SRC)

    # 1번째 인덱스 (Slide 2): 팀원 카드 등
    if len(prs.slides) > 1:
        rebuild_slide2(prs.slides[1])
    
    # 6번째 인덱스 (Slide 7): 테스트 케이스 표 (원래 6번째가 7번째 슬라이드)
    # 06. 테스트 케이스 및 결과 가 7번째 슬라이드일 것.
    if len(prs.slides) > 6:
        rebuild_slide7(prs.slides[6])
        
    # 7번째 인덱스 (Slide 8): 07 학습내용 & 교훈
    if len(prs.slides) > 7:
        rebuild_slide8(prs.slides[7])

    # 9번째 인덱스 (Slide 10)
    if len(prs.slides) > 9:
        fix_slide10(prs.slides[9])

    print(f'Saving: {OUT}')
    prs.save(OUT)
    print('✅ Done!')


if __name__ == '__main__':
    main()
