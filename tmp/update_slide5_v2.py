import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb(r, g, b):
    return RGBColor(r, g, b)

BLUE_DARK  = rgb(0x00, 0x1f, 0x54) # Dark navy blue
ORANGE     = rgb(0xff, 0x8c, 0x00)
WHITE      = rgb(0xff, 0xff, 0xff)
GRAY_BG    = rgb(0xf8, 0xf9, 0xfa)
DARK_TEXT  = rgb(0x21, 0x25, 0x29)
GREEN_SUCC = rgb(0x00, 0xc8, 0x53)
TERM_BG    = rgb(0x1e, 0x1e, 0x1e)
TERM_GREEN = rgb(0x81, 0xc7, 0x84)
TERM_WHITE = rgb(0xe0, 0xe0, 0xe0)
LITE_BLUE  = rgb(0xe3, 0xf2, 0xfd)
BORDER_GRAY= rgb(0xde, 0xe2, 0xe6)

def add_rect(slide, left, top, width, height, fill_color, border_color=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_tb(slide, left, top, width, height, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT, font_name='Malgun Gothic'):
    if color is None: color = DARK_TEXT
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txb

def clear_shapes(slide):
    # Remove all shapes except the title (but we'll clear the title's text so it doesn't overlap our header)
    to_remove = []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            if shape.has_text_frame:
                shape.text = ""
        else:
            to_remove.append(shape)
    
    for s in to_remove:
        s.element.getparent().remove(s.element)

def build_slide(slide):
    clear_shapes(slide)
    
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)

    # Top Header Bar
    add_rect(slide, 0, 0, 13.33, 0.8, BLUE_DARK)
    add_tb(slide, 0.4, 0.15, 12.0, 0.5, "05. 구현 결과 (Implementation Results)", size=24, bold=True, color=WHITE, font_name='Malgun Gothic')
    add_rect(slide, 0, 0.8, 13.33, 0.05, ORANGE)

    col_y = 1.2
    col_w = 6.2
    col_h = 5.8
    gap = 0.5
    total_w = col_w * 2 + gap
    start_x = (13.33 - total_w) / 2
    left_x = start_x
    right_x = start_x + col_w + gap

    # ---- Left Column Card ----
    add_rect(slide, left_x, col_y, col_w, col_h, WHITE, border_color=BORDER_GRAY, rounded=True)
    add_tb(slide, left_x, col_y+0.15, col_w, 0.4, "시뮬레이션 및 테스트 검증", size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # Add two line explanation
    desc_y = col_y + 0.55
    add_tb(slide, left_x + 0.2, desc_y, col_w - 0.4, 0.6, 
           "• GTest 프레임워크 기반 140개 단위/통합 테스트 100% 통과\n"
           "• 시뮬레이션 환경의 동적 입력(차속/센서)에 따른 정상 상태 전이 검증", 
           size=12, color=DARK_TEXT)

    # Left Column: Terminal window
    term_y = col_y + 1.2
    term_h = 1.9
    term_pad = 0.2
    add_rect(slide, left_x + term_pad, term_y, col_w - 2*term_pad, term_h, TERM_BG, rounded=True)
    
    dot_y = term_y + 0.1
    add_rect(slide, left_x + term_pad + 0.1, dot_y, 0.12, 0.12, rgb(0xff, 0x5f, 0x56), rounded=True)
    add_rect(slide, left_x + term_pad + 0.3, dot_y, 0.12, 0.12, rgb(0xff, 0xbd, 0x2e), rounded=True)
    add_rect(slide, left_x + term_pad + 0.5, dot_y, 0.12, 0.12, rgb(0x27, 0xc9, 0x3f), rounded=True)

    txb = slide.shapes.add_textbox(Inches(left_x + term_pad + 0.1), Inches(term_y + 0.3), Inches(col_w - 2*term_pad - 0.2), Inches(term_h - 0.4))
    tf = txb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "[==========] Running 140 tests from 12 test suites.\n"
    r.font.color.rgb = TERM_WHITE; r.font.name = 'Consolas'; r.font.size = Pt(11)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "[  PASSED  ] "
    r2.font.color.rgb = TERM_GREEN; r2.font.name = 'Consolas'; r2.font.size = Pt(11)
    r2b = p2.add_run()
    r2b.text = "140 tests.\n\n"
    r2b.font.color.rgb = TERM_WHITE; r2b.font.name = 'Consolas'; r2b.font.size = Pt(11)

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Coverage: 100% (Branch)\n"
    r3.font.color.rgb = TERM_WHITE; r3.font.name = 'Consolas'; r3.font.size = Pt(11)

    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    r4.text = "Cyclomatic Complexity Average: 7"
    r4.font.color.rgb = TERM_WHITE; r4.font.name = 'Consolas'; r4.font.size = Pt(11)

    # Left Column: Sequential log output snippet (Styled like a terminal as well)
    log_y = term_y + term_h + 0.2
    log_h = 2.0
    add_rect(slide, left_x + term_pad, log_y, col_w - 2*term_pad, log_h, TERM_BG, border_color=BORDER_GRAY, rounded=True) 
    
    log_txb = slide.shapes.add_textbox(Inches(left_x + term_pad + 0.1), Inches(log_y + 0.1), Inches(col_w - 2*term_pad - 0.2), Inches(log_h - 0.2))
    ltf = log_txb.text_frame
    ltf.word_wrap = True
    
    p0 = ltf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "=== Simulation System Log ===\n"
    r0.font.color.rgb = rgb(0xaa, 0xaa, 0xaa); r0.font.name = 'Consolas'; r0.font.size = Pt(10)

    def add_log_line(tf, prefix, transition):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = prefix + " "
        r.font.color.rgb = rgb(0x64, 0xb5, 0xf6) # light blue
        r.font.name = 'Consolas'; r.font.size = Pt(10)
        r2 = p.add_run()
        r2.text = transition
        r2.font.color.rgb = TERM_WHITE
        r2.font.name = 'Consolas'; r2.font.size = Pt(10)
    
    add_log_line(ltf, "[VehicleSpeed: 0km/h]", "UserRequest -> ManualLock")
    add_log_line(ltf, "[VehicleSpeed: 16km/h]", "-> AUTO_LOCKED")
    add_log_line(ltf, "[CollisionDetec: TRUE]", "-> EMERGENCY_RELEASED")


    # ---- Right Column Card ----
    add_rect(slide, right_x, col_y, col_w, col_h, WHITE, border_color=BORDER_GRAY, rounded=True)
    add_tb(slide, right_x, col_y+0.15, col_w, 0.4, "CI/CD 및 형상 관리", size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    desc_y2 = col_y + 0.55
    add_tb(slide, right_x + 0.2, desc_y2, col_w - 0.4, 0.6, 
           "• GitHub Actions 기반 빌드 ∙ 테스트 ∙ 정적 분석(Cppcheck) 자동화\n"
           "• 기능 개발(feat) 및 버그 수정(fix) 후 main 브랜치로 병합하는 워크플로우", 
           size=12, color=DARK_TEXT)

    # Right Column: GitHub Actions workflow
    gh_y = col_y + 1.2
    gh_h = 2.0
    gh_pad = 0.2
    add_rect(slide, right_x + gh_pad, gh_y, col_w - 2*gh_pad, gh_h, WHITE, border_color=BORDER_GRAY, rounded=True)
    
    add_tb(slide, right_x + gh_pad + 0.1, gh_y + 0.1, col_w - 0.4, 0.2, "GitHub Actions Pipeline", size=11, bold=True, color=DARK_TEXT)
    add_tb(slide, right_x + gh_pad + 0.1, gh_y + 0.35, col_w - 0.4, 0.2, "Commit: a1b2c3d  |  Duration: 1m 32s", size=9, color=rgb(0x6c, 0x75, 0x7d))
    
    jobs = [
        "Build (Ubuntu)",
        "Unit Tests (GTest)",
        "Static Analysis (Cppcheck)",
        "Code Complexity (Lizard)"
    ]
    job_y = gh_y + 0.6
    for job in jobs:
        add_tb(slide, right_x + gh_pad + 0.15, job_y - 0.05, 0.3, 0.3, "✔", size=12, bold=True, color=GREEN_SUCC, font_name='Segoe UI Symbol')
        add_tb(slide, right_x + gh_pad + 0.4, job_y - 0.02, 3.0, 0.25, job, size=10, bold=True, color=DARK_TEXT)
        job_y += 0.3
        
        # Add a subtle progress bar to make it look active
        add_rect(slide, right_x + gh_pad + 2.5, job_y - 0.22, 1.5, 0.08, LITE_BLUE, rounded=True)
        add_rect(slide, right_x + gh_pad + 2.5, job_y - 0.22, 1.5, 0.08, GREEN_SUCC, rounded=True)

    # Right Column: Git Branch Network Graph
    git_y = gh_y + gh_h + 0.2
    git_h = 1.9
    add_rect(slide, right_x + gh_pad, git_y, col_w - 2*gh_pad, git_h, WHITE, border_color=BORDER_GRAY, rounded=True)
    add_tb(slide, right_x + gh_pad + 0.1, git_y + 0.1, col_w - 0.4, 0.2, "Branch Network Graph", size=11, bold=True, color=DARK_TEXT)
    
    main_y = git_y + 1.0
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x + gh_pad + 0.3), Inches(main_y), Inches(5.0), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_DARK; shape.line.fill.background()
    add_tb(slide, right_x + gh_pad + 0.3, main_y - 0.3, 0.8, 0.2, "main", size=9, bold=True, color=BLUE_DARK)
    
    feat1_y = git_y + 0.5
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x + gh_pad + 0.8), Inches(feat1_y), Inches(1.7), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = ORANGE; shape.line.fill.background()
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x + gh_pad + 2.5), Inches(feat1_y), Inches(0.04), Inches(main_y - feat1_y))
    shape.fill.solid(); shape.fill.fore_color.rgb = ORANGE; shape.line.fill.background()
    
    add_tb(slide, right_x + gh_pad + 0.8, feat1_y - 0.25, 1.5, 0.2, "feat/rear-radar", size=9, color=ORANGE, bold=True)
    
    fix_y = git_y + 1.5
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x + gh_pad + 1.5), Inches(fix_y), Inches(1.5), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(0x00, 0xbc, 0xd4); shape.line.fill.background()
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x + gh_pad + 3.0), Inches(main_y), Inches(0.04), Inches(fix_y - main_y))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(0x00, 0xbc, 0xd4); shape.line.fill.background()
    
    add_tb(slide, right_x + gh_pad + 1.5, fix_y + 0.05, 1.5, 0.2, "fix/crash-logic", size=9, color=rgb(0x00, 0xbc, 0xd4), bold=True)
    
    nodes_x = [0.8, 1.5, 2.5, 3.0, 3.8, 4.5]
    for i, nx in enumerate(nodes_x):
        cx = right_x + gh_pad + nx
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.08), Inches(main_y - 0.055), Inches(0.16), Inches(0.16))
        circle.fill.solid(); circle.fill.fore_color.rgb = WHITE
        circle.line.color.rgb = BLUE_DARK; circle.line.width = Pt(2)
        
        if nx == 2.5: # merge from feat1
            c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.1), Inches(main_y - 0.075), Inches(0.2), Inches(0.2))
            c2.fill.solid(); c2.fill.fore_color.rgb = ORANGE; c2.line.fill.background()
            
        if nx == 3.0: # merge from fix
            c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.1), Inches(main_y - 0.075), Inches(0.2), Inches(0.2))
            c3.fill.solid(); c3.fill.fore_color.rgb = rgb(0x00, 0xbc, 0xd4); c3.line.fill.background()

        if nx == 4.5: # final merge PR
            mpr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.35), Inches(main_y - 0.35), Inches(0.7), Inches(0.22))
            mpr.fill.solid(); mpr.fill.fore_color.rgb = GREEN_SUCC; mpr.line.fill.background()
            add_tb(slide, cx - 0.35, main_y - 0.33, 0.7, 0.22, "Merge PR", size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer Watermark
    add_tb(slide, 11.0, 7.1, 2.0, 0.3, "HYUNDAI MOBIS", size=12, bold=True, color=rgb(0xce, 0xd4, 0xda), align=PP_ALIGN.RIGHT)

def main():
    path = r'c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\MBP-T05_ChildLock_결과발표_Final_v2.pptx'
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        return
        
    target_idx = -1
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                if '05. 구현 결과' in shape.text:
                    target_idx = i
                    break
        if target_idx != -1:
            break
            
    if target_idx == -1:
        print("Could not find Slide with '05. 구현 결과'. Applying to Slide 5 (index 4).")
        target_idx = 4
        
    print(f"Applying changes to slide index {target_idx}")
    build_slide(prs.slides[target_idx])
    
    prs.save(path)
    print("Done!")

if __name__ == '__main__':
    main()
