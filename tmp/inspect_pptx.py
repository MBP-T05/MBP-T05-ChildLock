from pptx import Presentation
import sys

with open('tmp/pptx_layout.txt', 'w', encoding='utf-8') as f:
    prs = Presentation('docs/MOBIUS_PBL_결과발표_템플릿.pptx')
    f.write("--- Slide Layouts ---\n")
    for i, layout in enumerate(prs.slide_layouts):
        f.write(f"[{i}] {layout.name}\n")
        for j, shape in enumerate(layout.placeholders):
            f.write(f"  - Pld {shape.placeholder_format.idx}: {shape.name}\n")

    f.write("\n--- Slides Text ---\n")
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                f.write(shape.text + "\n")
