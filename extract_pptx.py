import collections
import collections.abc
import sys
try:
    import pptx
except ImportError:
    pass

def extract_text(pptx_path, out_path):
    with open(out_path, 'w', encoding='utf-8') as f:
        try:
            prs = pptx.Presentation(pptx_path)
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        f.write(shape.text.strip() + "\n")
                f.write("\n")
            print("Extracted PPTX successfully.")
        except Exception as e:
            f.write(f"Error opening PPTX: {e}\n")
            print(f"Error opening PPTX: {e}")

if __name__ == "__main__":
    if len(sys.argv) == 3:
        extract_text(sys.argv[1], sys.argv[2])
    else:
        extract_text(r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\MBP-T05_ChildLock_결과발표.pptx", r"C:\Users\fish0\.gemini\antigravity\brain\d696c5d4-907c-4e6e-bc04-2a070b2afe2f\pptx_content.txt")
